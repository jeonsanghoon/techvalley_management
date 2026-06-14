import json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

META_FILE = "/tmp/tv-screenshots/meta.json"
OUT_FILE  = "/Users/jeonsanghun/Downloads/techvalley-screens.pptx"

SLIDE_W = Inches(20)   # 1920/96 dpi
SLIDE_H = Inches(11.25) # 1080/96 dpi

TITLE_H  = Inches(0.45)
IMG_TOP  = TITLE_H
IMG_H    = SLIDE_H - TITLE_H

DARK_BG  = RGBColor(0x0D, 0x0D, 0x0F)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)

def add_slide(prs, img_path, title_text, mode):
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)

    # 배경
    bg_color = DARK_BG if mode == "dark" else LIGHT_BG
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # 타이틀 바
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(0), Emu(0), SLIDE_W, TITLE_H,
    )
    bar.fill.solid()
    bar_color = RGBColor(0x1A, 0x1E, 0x24) if mode == "dark" else RGBColor(0x1E, 0x40, 0x80)
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"  {title_text}"
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.size = Pt(16)
    run.font.bold = True

    # 모드 뱃지
    badge_w = Inches(1.2)
    badge = slide.shapes.add_shape(
        1, SLIDE_W - badge_w, Emu(0), badge_w, TITLE_H
    )
    badge.fill.solid()
    badge_color = RGBColor(0x2C, 0xB6, 0x7D) if mode == "dark" else RGBColor(0x3B, 0x82, 0xF6)
    badge.fill.fore_color.rgb = badge_color
    badge.line.fill.background()
    bt = badge.text_frame.paragraphs[0]
    bt.alignment = PP_ALIGN.CENTER
    br = bt.add_run()
    br.text = "🌙 Dark" if mode == "dark" else "☀️ Light"
    br.font.size = Pt(13)
    br.font.bold = True
    br.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 스크린샷
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Emu(0), IMG_TOP, SLIDE_W, IMG_H)

    return slide


def main():
    with open(META_FILE) as f:
        meta = json.load(f)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # 표지
    cover_layout = prs.slide_layouts[6]
    cover = prs.slides.add_slide(cover_layout)
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    title_box = cover.shapes.add_textbox(Inches(2), Inches(4), Inches(16), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "테크밸리 IoT 서비스 플랫폼"
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    sub_box = cover.shapes.add_textbox(Inches(2), Inches(6), Inches(16), Inches(1))
    sp = sub_box.text_frame.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = "화면 구성 — 라이트 모드 / 다크 모드"
    sr.font.size = Pt(22)
    sr.font.color.rgb = RGBColor(0xA0, 0xB8, 0xD8)

    # 구분 슬라이드 — 라이트
    def section_slide(label, color):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = color
        tb = sl.shapes.add_textbox(Inches(2), Inches(4.5), Inches(16), Inches(2))
        p2 = tb.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(48)
        r2.font.bold = True
        r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    section_slide("☀️  라이트 모드", RGBColor(0x1E, 0x40, 0x80))

    for item in meta.get("light", []):
        add_slide(prs, item["file"], item["title"], "light")

    section_slide("🌙  다크 모드", RGBColor(0x1A, 0x1E, 0x24))

    for item in meta.get("dark", []):
        add_slide(prs, item["file"], item["title"], "dark")

    prs.save(OUT_FILE)
    total = len(meta.get("light", [])) + len(meta.get("dark", []))
    print(f"✅ PPT 저장 완료: {OUT_FILE}  ({total}슬라이드 + 표지/구분 3장)")


if __name__ == "__main__":
    main()
