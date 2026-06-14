#!/usr/bin/env python3
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W  = Inches(20)
SLIDE_H  = Inches(11.25)
META_FILE= "/tmp/tv-screenshots/meta.json"
OUT_FILE = "/Users/jeonsanghun/Downloads/techvalley-manual.pptx"

# ── 색상 ──────────────────────────────────────────────────────────
CB  = RGBColor(0x0F,0x17,0x2A)   # bg dark
CC  = RGBColor(0x16,0x22,0x38)   # card
CH  = RGBColor(0x1A,0x2B,0x4A)   # header
BLU = RGBColor(0x3B,0x82,0xF6)
CYN = RGBColor(0x06,0xB6,0xD4)
GRN = RGBColor(0x10,0xB9,0x81)
AMB = RGBColor(0xF5,0x9E,0x0B)
RED = RGBColor(0xEF,0x44,0x44)
PRP = RGBColor(0x8B,0x5C,0xF6)
TEA = RGBColor(0x14,0xB8,0xA6)
IND = RGBColor(0x63,0x66,0xF1)
WHT = RGBColor(0xFF,0xFF,0xFF)
GRY = RGBColor(0x94,0xA3,0xB8)
LGT = RGBColor(0xCB,0xD5,0xE1)

# ── 헬퍼 ─────────────────────────────────────────────────────────
def prs_new():
    p = Presentation(); p.slide_width=SLIDE_W; p.slide_height=SLIDE_H; return p

def blank(prs, bg=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if bg: s.background.fill.solid(); s.background.fill.fore_color.rgb=bg
    return s

def R(s,l,t,w,h,fill,line=None,rnd=False):
    sh=s.shapes.add_shape(5 if rnd else 1,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(1)
    else: sh.line.fill.background()
    return sh

def T(s,text,l,t,w,h,sz=16,bold=False,col=WHT,align=PP_ALIGN.LEFT,wrap=True):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=text; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col
    return tb

def add_ss(s,path,l,t,w):
    if path and os.path.exists(path):
        s.shapes.add_picture(path,l,t,w,w*9/16)

# ── 스크린샷 맵 ──────────────────────────────────────────────────
with open(META_FILE) as f: meta=json.load(f)
SS={}       # path → first screenshot
SS_TABS={}  # path → [{tab, file, title}, ...]
for it in meta["light"]:
    p=it["path"]
    if p not in SS: SS[p]=it["file"]
    if it.get("tab"):
        SS_TABS.setdefault(p,[]).append(it)

# ══════════════════════════════════════════════════════════════════
# 슬라이드 빌더
# ══════════════════════════════════════════════════════════════════
def cover(prs):
    s=blank(prs,CB)
    R(s,0,0,SLIDE_W,Inches(0.08),BLU)
    R(s,0,SLIDE_H-Inches(0.08),SLIDE_W,Inches(0.08),BLU)
    T(s,"TECHVALLEY",Inches(1),Inches(1.2),Inches(18),Inches(1.8),sz=68,bold=True,col=BLU,align=PP_ALIGN.CENTER)
    T(s,"IoT 서비스 플랫폼",Inches(1),Inches(3.1),Inches(18),Inches(1.0),sz=34,col=LGT,align=PP_ALIGN.CENTER)
    R(s,Inches(5),Inches(4.4),Inches(10),Inches(0.05),BLU)
    T(s,"사 용 자  설 명 서",Inches(1),Inches(4.7),Inches(18),Inches(1.8),sz=54,bold=True,col=WHT,align=PP_ALIGN.CENTER)
    T(s,"User Operation Guide",Inches(1),Inches(6.6),Inches(18),Inches(0.8),sz=22,col=GRY,align=PP_ALIGN.CENTER)
    chips=["실시간 모니터링","알람 관리","서비스 티켓","원격 제어","부품 조달","보고서"]
    cw,cg=Inches(2.2),Inches(0.3); tot=len(chips)*cw+(len(chips)-1)*cg
    sx=(SLIDE_W-tot)/2
    for i,c in enumerate(chips):
        cx=sx+i*(cw+cg)
        R(s,cx,Inches(8.0),cw,Inches(0.55),CH,rnd=True)
        T(s,c,cx,Inches(8.03),cw,Inches(0.5),sz=14,col=CYN,align=PP_ALIGN.CENTER)
    T(s,"Ver 1.0  ·  2026",Inches(1),Inches(9.8),Inches(18),Inches(0.6),sz=15,col=GRY,align=PP_ALIGN.CENTER)

def toc(prs):
    s=blank(prs,CB)
    R(s,0,0,SLIDE_W,Inches(1.1),CH)
    T(s,"목  차",Inches(0.5),Inches(0.2),Inches(10),Inches(0.7),sz=30,bold=True)
    T(s,"Table of Contents",Inches(0.5),Inches(0.7),Inches(10),Inches(0.4),sz=15,col=GRY)
    rows=[
        ("01","시스템 개요","플랫폼 구성 및 주요 기능 소개"),
        ("02","장비 관리","장비 라이프사이클 · 목록 · 이력"),
        ("03","데이터 수집 / 모니터링","파이프라인 · 메트릭 스트림"),
        ("04","알람 관리","알람 처리 흐름 · 규칙 설정"),
        ("05","서비스 관리","서비스 티켓 · 진행 · A/S"),
        ("06","현장 작업","점검 · 설치 · 원격 제어"),
        ("07","부품 관리","부품 조달 흐름 · 주문 · 일정"),
        ("08","고객 / SLA 관리","고객사 · 서비스 레벨"),
        ("09","보고서 / 설정","보고서 · 알림 · 펌웨어 · IoT 인증"),
        ("10","통합 대시보드","실시간 KPI · 현황판"),
    ]
    cw=Inches(9.0); gap=Inches(0.5)
    for i,(num,title,sub) in enumerate(rows):
        col=i%2; row=i//2
        lx=Inches(0.5)+col*(cw+gap); ly=Inches(1.3)+row*Inches(1.88)
        R(s,lx,ly+Inches(0.1),Inches(0.75),Inches(0.75),BLU,rnd=True)
        T(s,num,lx,ly+Inches(0.1),Inches(0.75),Inches(0.75),sz=19,bold=True,align=PP_ALIGN.CENTER)
        R(s,lx+Inches(0.9),ly,cw-Inches(0.9),Inches(0.95),CC,rnd=True)
        T(s,title,lx+Inches(1.05),ly+Inches(0.05),cw-Inches(1.1),Inches(0.5),sz=19,bold=True)
        T(s,sub,lx+Inches(1.05),ly+Inches(0.55),cw-Inches(1.1),Inches(0.38),sz=13,col=GRY)

def overview(prs):
    s=blank(prs,CB)
    R(s,0,0,SLIDE_W,Inches(1.05),CH)
    T(s,"시스템 개요",Inches(0.5),Inches(0.2),Inches(14),Inches(0.65),sz=28,bold=True)
    T(s,"System Overview",Inches(0.5),Inches(0.7),Inches(14),Inches(0.38),sz=14,col=GRY)
    T(s,"테크밸리 IoT 서비스 플랫폼은 현장에 배포된 의료/산업용 장비의 실시간 모니터링, 알람 처리, 현장 서비스 관리를 통합 제공합니다.",
       Inches(0.5),Inches(1.2),Inches(19),Inches(0.65),sz=16,col=LGT)
    blocks=[
        ("📡","실시간\n모니터링","장비 센서 데이터 수집\n실시간 분석·가시화",BLU),
        ("🔔","알람\n관리","임계값 기반 자동 감지\n담당자 즉시 알림",AMB),
        ("🛠","서비스\n티켓","접수→배정→출동→완료\n전 단계 추적",PRP),
        ("🎮","원격\n제어","파라미터 원격 조정\n안전 모드 제어",CYN),
        ("📦","부품\n조달","재고→발주→입고→교체\n공급망 관리",GRN),
        ("📊","통합\n대시보드","KPI 실시간 현황\n플릿 전체 가시성",RGBColor(0xE8,0x78,0x3A)),
    ]
    bw=Inches(2.8); bg=Inches(0.2); tot=len(blocks)*bw+(len(blocks)-1)*bg
    bx=(SLIDE_W-tot)/2; by=Inches(2.1); bh=Inches(5.8)
    for i,(ic,ti,de,co) in enumerate(blocks):
        x=bx+i*(bw+bg)
        R(s,x,by,bw,bh,CC,rnd=True)
        R(s,x,by,bw,Inches(0.1),co)
        T(s,ic,x,by+Inches(0.25),bw,Inches(0.9),sz=36,align=PP_ALIGN.CENTER)
        T(s,ti,x,by+Inches(1.2),bw,Inches(0.9),sz=17,bold=True,align=PP_ALIGN.CENTER)
        T(s,de,x+Inches(0.1),by+Inches(2.2),bw-Inches(0.2),Inches(1.5),sz=13,col=LGT,align=PP_ALIGN.CENTER)
    T(s,"현장 장비  →  IoT 게이트웨이  →  클라우드 처리  →  서비스 플랫폼  →  담당자 / 관리자",
       Inches(0.5),Inches(8.3),Inches(19),Inches(0.7),sz=16,col=CYN,align=PP_ALIGN.CENTER)
    R(s,Inches(0.5),Inches(8.25),Inches(19),Inches(0.04),CH)

def sec_slide(prs,num,title,sub,ac):
    s=blank(prs,CB)
    R(s,0,0,Inches(0.15),SLIDE_H,ac)
    T(s,num,Inches(11),Inches(0.5),Inches(9),Inches(9),sz=280,bold=True,
      col=RGBColor(0x18,0x25,0x3A),align=PP_ALIGN.RIGHT)
    T(s,f"SECTION {num}",Inches(0.5),Inches(3.2),Inches(14),Inches(0.7),sz=18,col=GRY)
    T(s,title,Inches(0.5),Inches(4.0),Inches(14),Inches(1.5),sz=52,bold=True)
    R(s,Inches(0.5),Inches(5.7),Inches(4),Inches(0.06),ac)
    T(s,sub,Inches(0.5),Inches(5.9),Inches(14),Inches(0.9),sz=20,col=LGT)

def lifecycle(prs,title,steps,note=None):
    s=blank(prs,CB)
    R(s,0,0,SLIDE_W,Inches(1.05),CH)
    T(s,f"🔄  {title}",Inches(0.5),Inches(0.2),Inches(16),Inches(0.65),sz=26,bold=True)
    R(s,Inches(17.5),Inches(0.22),Inches(2.1),Inches(0.58),BLU,rnd=True)
    T(s,"LIFECYCLE",Inches(17.5),Inches(0.26),Inches(2.1),Inches(0.5),sz=13,bold=True,align=PP_ALIGN.CENTER)

    n=len(steps); SW=Inches(2.4); SH=Inches(4.5); AW=Inches(0.4); AH=Inches(0.35)
    tot=n*SW+(n-1)*AW; sx=(SLIDE_W-tot)/2; sy=Inches(2.85)
    for i,st in enumerate(steps):
        x=sx+i*(SW+AW)
        bx=s.shapes.add_shape(5,x,sy,SW,SH)
        bx.fill.solid(); bx.fill.fore_color.rgb=st["color"]; bx.line.fill.background()
        tf=bx.text_frame; tf.word_wrap=True
        p0=tf.paragraphs[0]; p0.alignment=PP_ALIGN.CENTER; p0.space_before=Pt(18)
        r0=p0.add_run(); r0.text=str(i+1); r0.font.size=Pt(36); r0.font.bold=True; r0.font.color.rgb=WHT
        p1=tf.add_paragraph(); p1.alignment=PP_ALIGN.CENTER; p1.space_before=Pt(6)
        r1=p1.add_run(); r1.text=st["title"]; r1.font.size=Pt(16); r1.font.bold=True; r1.font.color.rgb=WHT
        if st.get("desc"):
            p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER; p2.space_before=Pt(10)
            r2=p2.add_run(); r2.text=st["desc"]; r2.font.size=Pt(12); r2.font.color.rgb=LGT
        if i<n-1:
            ax=x+SW; ay=sy+(SH-AH)/2
            ar=s.shapes.add_shape(13,ax,ay,AW,AH)
            ar.fill.solid(); ar.fill.fore_color.rgb=GRY; ar.line.fill.background()
    if note:
        R(s,Inches(0.5),Inches(9.1),Inches(19),Inches(0.04),CC)
        T(s,f"💡  {note}",Inches(0.5),Inches(9.2),Inches(19),Inches(0.7),sz=14,col=GRY)

def feat(prs,title,icon,desc,features,actions,path,ac=None):
    ac=ac or BLU
    s=blank(prs,RGBColor(0x0D,0x14,0x23))
    # 좌측 패널
    R(s,0,0,Inches(7.6),SLIDE_H,CC)
    R(s,0,0,Inches(0.12),SLIDE_H,ac)
    T(s,icon,Inches(0.25),Inches(0.35),Inches(1.1),Inches(1.0),sz=42,align=PP_ALIGN.CENTER)
    T(s,title,Inches(1.4),Inches(0.45),Inches(6.0),Inches(0.85),sz=26,bold=True)
    R(s,Inches(0.3),Inches(1.45),Inches(7.0),Inches(0.04),ac)
    # 설명
    tb=s.shapes.add_textbox(Inches(0.3),Inches(1.6),Inches(7.0),Inches(2.1))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=desc; r.font.size=Pt(14); r.font.color.rgb=LGT
    # 주요 기능
    R(s,Inches(0.3),Inches(3.85),Inches(0.08),Inches(0.5),ac)
    T(s,"주요 기능",Inches(0.5),Inches(3.85),Inches(3),Inches(0.48),sz=15,bold=True,col=ac)
    ftb=s.shapes.add_textbox(Inches(0.3),Inches(4.4),Inches(7.1),Inches(3.5))
    ftf=ftb.text_frame; ftf.word_wrap=True
    for j,f in enumerate(features):
        p2=ftf.paragraphs[0] if j==0 else ftf.add_paragraph()
        p2.space_before=Pt(5); r2=p2.add_run(); r2.text=f"  •  {f}"; r2.font.size=Pt(13); r2.font.color.rgb=LGT
    # 주요 액션
    ab=Inches(4.4+len(features)*0.42+0.3)
    R(s,Inches(0.3),ab,Inches(0.08),Inches(0.5),AMB)
    T(s,"주요 액션",Inches(0.5),ab,Inches(3),Inches(0.48),sz=15,bold=True,col=AMB)
    atb=s.shapes.add_textbox(Inches(0.3),ab+Inches(0.55),Inches(7.1),Inches(2.0))
    atf=atb.text_frame; atf.word_wrap=True
    for j,a in enumerate(actions):
        p3=atf.paragraphs[0] if j==0 else atf.add_paragraph()
        p3.space_before=Pt(5); r3=p3.add_run(); r3.text=f"  ▶  {a}"; r3.font.size=Pt(13); r3.font.color.rgb=LGT
    # 우측 스크린샷
    rx=Inches(7.9); rw=SLIDE_W-rx-Inches(0.2)
    ss=SS.get(path)
    if ss and os.path.exists(ss):
        sh=rw*9/16
        ry=(SLIDE_H-sh)/2
        s.shapes.add_picture(ss,rx,ry,rw,sh)
    else:
        R(s,rx,Inches(0.3),rw,SLIDE_H-Inches(0.6),RGBColor(0x1A,0x24,0x38))

def tab_grid(prs, page_title, icon, tabs_info, ac=None):
    """탭 스크린샷을 그리드(2×N)로 보여주는 슬라이드들 추가"""
    ac = ac or CYN
    # 한 슬라이드에 최대 4개 탭 (2×2)
    chunk = 4
    for ci in range(0, len(tabs_info), chunk):
        group = tabs_info[ci:ci+chunk]
        s = blank(prs, RGBColor(0x0D,0x14,0x23))
        R(s,0,0,SLIDE_W,Inches(0.85),CH)
        T(s,f"{icon}  {page_title}  —  탭별 화면",Inches(0.4),Inches(0.15),Inches(16),Inches(0.58),sz=22,bold=True)
        badge_txt = f"{ci+1}–{min(ci+chunk,len(tabs_info))} / {len(tabs_info)} 탭"
        T(s,badge_txt,Inches(16.5),Inches(0.2),Inches(3.2),Inches(0.5),sz=14,col=GRY,align=PP_ALIGN.RIGHT)

        cols=2; rows=2
        padx=Inches(0.4); pady=Inches(0.15)
        cw=(SLIDE_W-padx*2)/cols; ch=(SLIDE_H-Inches(0.85)-pady*2)/rows
        img_h=cw*9/16

        for gi,it in enumerate(group):
            col=gi%cols; row=gi//cols
            lx=padx+col*cw; ty=Inches(0.85)+pady+row*ch
            # 탭 레이블 바
            R(s,lx+Inches(0.05),ty+Inches(0.05),cw-Inches(0.1),Inches(0.4),ac,rnd=True)
            T(s,it["title"].split("—")[-1].strip(),
              lx+Inches(0.1),ty+Inches(0.05),cw-Inches(0.2),Inches(0.4),
              sz=15,bold=True,align=PP_ALIGN.CENTER)
            # 스크린샷
            f=it["file"]
            if os.path.exists(f):
                iy=ty+Inches(0.5)
                s.shapes.add_picture(f,lx+Inches(0.05),iy,cw-Inches(0.1),img_h)

# ══════════════════════════════════════════════════════════════════
# 라이프사이클 데이터
# ══════════════════════════════════════════════════════════════════
LC_EQUIP={
    "title":"장비 관리 라이프사이클",
    "steps":[
        {"title":"고객사 등록","desc":"고객사 정보\n사이트 등록","color":BLU},
        {"title":"설치 계획","desc":"설치 일정\n엔지니어 배정","color":RGBColor(0x2D,0x6E,0xE0)},
        {"title":"IoT 연동","desc":"Thing 등록\n인증서 발급","color":CYN},
        {"title":"정상 운영","desc":"데이터 수집\n실시간 모니터링","color":GRN},
        {"title":"이상 감지","desc":"알람 발생\n티켓 생성","color":AMB},
        {"title":"서비스 처리","desc":"엔지니어 출동\n수리/교체 후 복귀","color":RED},
    ],
    "note":"서비스 완료 후 '정상 운영' 단계로 복귀하여 지속적인 모니터링을 수행합니다.",
}
LC_ALARM={
    "title":"알람 처리 라이프사이클",
    "steps":[
        {"title":"규칙 설정","desc":"임계값/조건\n규칙 등록","color":IND},
        {"title":"임계값 감지","desc":"센서 데이터\n실시간 분석","color":PRP},
        {"title":"알람 발생","desc":"Critical/Warning\n알림 전송","color":RED},
        {"title":"담당자 확인","desc":"알람 확인\n원인 분석","color":AMB},
        {"title":"티켓 생성","desc":"서비스 티켓\n자동/수동 생성","color":RGBColor(0xF5,0x7A,0x22)},
        {"title":"알람 해제","desc":"처리 완료\n이력 기록","color":GRN},
    ],
    "note":"Critical 알람은 즉시 SNS/이메일 알림이 발송되며, Unacked 상태로 미확인 알람이 관리됩니다.",
}
LC_SVC={
    "title":"서비스 티켓 라이프사이클",
    "steps":[
        {"title":"접  수","desc":"Received\n티켓 등록","color":BLU},
        {"title":"배  정","desc":"Assigned\n엔지니어 지정","color":RGBColor(0x48,0x40,0xC0)},
        {"title":"출  동","desc":"Dispatched\n현장 이동","color":PRP},
        {"title":"작업 중","desc":"Working\n현장 처리","color":AMB},
        {"title":"완  료","desc":"Completed\n처리 결과 입력","color":GRN},
        {"title":"A/S 등록","desc":"만족도 평가\n이력 기록","color":TEA},
    ],
    "note":"SLA 기준 초과 티켓은 'SLA At Risk' 상태로 표시되어 우선 처리됩니다.",
}
LC_PARTS={
    "title":"부품 조달 라이프사이클",
    "steps":[
        {"title":"주문 요청","desc":"Request\n필요 부품 발주","color":BLU},
        {"title":"주문 확인","desc":"Confirmed\n공급사 확인","color":CYN},
        {"title":"출  하","desc":"Shipped\n창고 출고","color":RGBColor(0x10,0x9A,0x78)},
        {"title":"운송 중","desc":"In Transit\n배송 추적","color":AMB},
        {"title":"도  착","desc":"Arrived\n입고 확인","color":RGBColor(0xEA,0x7A,0x1E)},
        {"title":"교체 완료","desc":"Replaced\nPOD 확인","color":GRN},
    ],
    "note":"부품 일정 화면에서 ETA · 방문일 · 엔지니어 배정 · POD 상태를 통합 관리합니다.",
}
LC_FIELD={
    "title":"현장 작업 흐름 (설치 · 점검)",
    "steps":[
        {"title":"설치 계획","desc":"일정 등록\n장비·사이트 지정","color":BLU},
        {"title":"현장 도착","desc":"엔지니어 출동\n장비 확인","color":RGBColor(0x28,0x6E,0xD0)},
        {"title":"IoT 등록","desc":"Thing 인증\n네트워크 연결","color":CYN},
        {"title":"시운전","desc":"파라미터 설정\n동작 확인","color":AMB},
        {"title":"알고리즘\n검사","desc":"수율 측정\n임계값 적용","color":PRP},
        {"title":"완료 확인","desc":"Commissioning\n고객 인수","color":GRN},
    ],
    "note":"설치 완료 후 IoT 자동 수집이 시작되며, 데이터 파이프라인에 장비가 자동 등록됩니다.",
}

# ══════════════════════════════════════════════════════════════════
# 화면별 기능 설명 데이터
# ══════════════════════════════════════════════════════════════════
PAGES=[
  ("/dashboard","통합 대시보드","📊",
   "현장 배포 장비의 운영 현황을 실시간 파악하는 중앙 관제 화면입니다. KPI 지표, 알람 현황, SLA 성과, 장비별 상태를 한눈에 확인합니다.",
   ["KPI 스트립: 가동률·알람·유지보수·티켓·수율 지표","플릿 지도: 지역별 장비 위치 및 상태","최근 알람 목록 및 빠른 링크","차트: 플릿 상태 도넛·티켓 단계·알람 트렌드·수율 게이지","장비 상태 그리드 및 활성 티켓 그리드"],
   ["알람 전체 보기","장비 목록으로 이동","서비스 티켓으로 이동"],BLU),

  ("/equipment","장비 목록","🏭",
   "현장 배포 전체 장비 인벤토리를 관리합니다. 지역·상태별 필터링으로 원하는 장비를 빠르게 찾고 상태를 확인합니다.",
   ["통계 그리드: 전체·온라인·알람·수명 낮음 수량","테이블: S/N·고객사·사이트·모델·상태·지역","상태 필터: Online / Alarm / Maintenance / Safe Mode / Offline","지역 필터 및 통합 검색"],
   ["신규 장비 등록","장비 이력 보기","원격 제어로 이동"],CYN),

  ("/equipment-logs","장비 이력","📋",
   "장비별 이벤트 이력을 카테고리·날짜 범위 별로 조회합니다. INFO/WARN/ERROR 레벨 로그를 통해 상태 변화를 추적합니다.",
   ["카테고리 탭: Tube·Detector·Body 등 부위별 분류","로그 레벨: INFO / WARN / ERROR 칩 필터","날짜 범위 프리셋 및 Warm Tier 배치 조회","메시지·소스·페이로드 검색"],
   ["실시간 메트릭 보기","날짜 범위 프리셋 적용","카테고리 전환"],CYN),

  ("/data-pipeline","데이터 파이프라인","🔌",
   "IoT 장비에서 수집되는 데이터의 저장 티어 상태(Hot/Warm/Cold)와 실시간 수집 현황을 모니터링합니다.",
   ["통계: 등록 장비 수·온라인 집계·당일 정규화·GG 컴포넌트","파이프라인 티어 테이블: Hot/Warm/Cold 상태","실시간 장비 수신 상태 목록","장비 S/N·고객사·사이트 검색 및 티어 필터"],
   ["상태별 필터링","티어별 조회","장비 수신 현황 확인"],BLU),

  ("/alarms","알람 목록","🔔",
   "발생 알람 이력을 심각도·확인 상태별로 관리합니다. Critical 미확인 알람은 상단 배너로 즉시 표시됩니다.",
   ["Critical 배너: 미확인 긴급 알람 즉시 표시","빠른 필터: 전체·미확인 Critical·미확인","통계: 전체·Critical·미확인·티켓 연결 수","날짜 범위·심각도·확인 상태 필터"],
   ["알람 일괄 확인","서비스 티켓 생성","원격 제어로 이동"],AMB),

  ("/alarm-rules","알람 규칙","⚙️",
   "알람 발생 조건과 심각도를 정의하는 규칙을 설정·관리합니다. 규칙별 활성화/비활성화 제어가 가능합니다.",
   ["규칙셋 테이블: 이름·대상·조건·심각도·활성화 상태","심각도 필터: Critical / Warning","활성화 토글 필터","규칙명·대상·조건 통합 검색"],
   ["규칙 추가","알림 채널 보기","활성화/비활성화 전환"],AMB),

  ("/service-tickets","서비스 티켓","🎫",
   "현장 서비스 요청 전체 목록을 관리합니다. Received → Assigned → Dispatched → Working → Completed 진행 상황을 추적합니다.",
   ["통계: 진행 중·Critical·SLA 위험·완료 수","단계별 필터: Received/Assigned/Dispatched/Working/Completed","티켓 ID·장비 S/N·고객사·증상 검색"],
   ["신규 티켓 발행","서비스 진행 화면으로 이동","단계별 필터링"],PRP),

  ("/service-progress","서비스 진행","🗓",
   "활성 서비스 티켓의 진행 상황을 엔지니어 가용성 그리드와 단계 추적 카드로 시각화합니다.",
   ["엔지니어 가용성: 이름·지역·상태·배정 수","단계 추적 카드: 티켓별 현재 단계 표시","티켓 상세 테이블","티켓 ID·장비·고객사·엔지니어 검색"],
   ["단계별 필터링","엔지니어 현황 확인"],PRP),

  ("/as","A/S 기록","🔖",
   "완료된 서비스 방문 이력을 관리합니다. 만족도 평가와 재발 여부를 추적하여 서비스 품질을 개선합니다.",
   ["통계: 기록 수·평균 만족도·재발 건수","유지보수 이력 테이블","재발 필터: 전체 / 재발 / 정상","티켓 ID·S/N·고객사·작업내용·교체부품 검색"],
   ["재발 필터링","이력 검색"],TEA),

  ("/inspection","점  검","🔬",
   "알고리즘 구성 및 수율 추적성을 관리합니다. 활성 알고리즘의 임계값 설정과 LOT별 수율 이력을 확인합니다.",
   ["통계: 평균 수율%·활성 알고리즘 버전·임계값·LOT 수","알고리즘 설정 테이블: Active/Staging/Disabled","현재 결정 기준 카드","수율/LOT 추적성 테이블"],
   ["알고리즘 등록","상태별 필터링"],RGBColor(0x6D,0x28,0xD9)),

  ("/installation","설  치","🔧",
   "장비 설치 프로젝트를 추적합니다. IoT 등록 상태와 시운전 진행도를 관리합니다.",
   ["통계: 전체·진행 중·IoT 등록 완료 수","설치/시운전 테이블: 발주참조·S/N·모델·고객사·사이트","상태 필터: Planned/In Progress/Commissioning/Completed","IoT 등록 상태 토글"],
   ["설치 일정 등록","상태별 필터링"],RGBColor(0x0F,0x76,0x6E)),

  ("/remote-control","원격 제어","🎮",
   "선택한 장비의 파라미터를 원격으로 조정합니다. kV/mA 슬라이더로 실시간 제어하며 안전 모드 및 긴급 연락 기능을 제공합니다.",
   ["장비 선택: S/N·모델·고객사·사이트 검색","파라미터 제어: kV (100~200) · mA (1~6) 슬라이더","안전 모드 전환 버튼 / EMG 호출","해상도 판정 결과 표시"],
   ["IoT 작업 실행 (파라미터 적용)","안전 모드 전환","EMG 호출","서비스 티켓 발행"],CYN),

  ("/parts-orders","부품 주문","📦",
   "교체 부품의 발주부터 교체 완료까지 공급망 전체를 추적합니다.",
   ["통계: 전체·운송 중·완료 수","주문 테이블: 주문ID·티켓ID·S/N·부품번호·부품명","상태: Request/Confirmed/Shipped/In Transit/Arrived/Replaced","통합 검색 지원"],
   ["주문 요청 생성","상태별 필터링"],GRN),

  ("/parts-schedule","부품 일정","🗓",
   "부품 배송 일정과 엔지니어 방문 계획을 통합 관리합니다. ETA·방문일·POD 상태를 추적합니다.",
   ["통계: 전체·운송 중·지연·완료 수","오늘 방문 예정 엔지니어 목록","일정 테이블: ETA·방문일·엔지니어·POD 상태","POD 상태 필터 및 통합 검색"],
   ["POD 상태별 필터링","배송 추적"],GRN),

  ("/customers","고객사 관리","🏢",
   "고객사와 딜러의 정보 및 설치 사이트를 관리합니다. 지역별·등록일별 필터링으로 계정 현황을 파악합니다.",
   ["통계: 고객사·사이트·총 고객·총 사이트 수","고객사/사이트 테이블 (주소·지역 포함)","유형 필터: Customer / Dealer","지역 필터 및 등록/설치일 범위 필터"],
   ["유형별 필터링","지역별 필터링"],RGBColor(0x0F,0x76,0x6E)),

  ("/sla","SLA 관리","📏",
   "서비스 레벨 계약(SLA) 티어 매트릭스와 장비별 가용성 현황을 관리합니다.",
   ["SLA 티어 테이블: 응답·출동·복구 시간·가동률%","장비 SLA 테이블: 가용성 지표","SLA 티어 필터: Critical / High / Standard","장비 S/N·고객사·서비스 윈도우 검색"],
   ["SLA 티어별 필터링"],RGBColor(0xDB,0x27,0x77)),

  ("/reports","보고서","📑",
   "운영·알람·서비스·점검 카테고리별 보고서를 생성하고 관리합니다.",
   ["카탈로그 테이블: 보고서명·카테고리·ID·생성일","카테고리 필터: Operations/Alarms/Service/Inspection","보고서명·카테고리·ID 검색"],
   ["보고서 생성"],RGBColor(0xDB,0x27,0x77)),

  ("/settings/notifications","알림 채널 설정","📬",
   "알람 발생 시 알림을 전송할 채널을 설정합니다. SNS, SES, 대시보드, Webhook 등 다양한 채널을 구성할 수 있습니다.",
   ["채널 테이블: 유형·대상·수신자·활성화 상태","유형 필터: SNS / SES / Dashboard / Webhook","활성화 토글 필터","채널명·대상·수신자 검색"],
   ["채널 추가","활성화/비활성화"],AMB),

  ("/settings/firmware","펌웨어 OTA 설정","💾",
   "장비 펌웨어 버전 추적과 OTA 업데이트 정책을 관리합니다. 자동 업데이트 설정 및 배포 작업을 실행합니다.",
   ["OTA 정책: 배포 스케줄 (02:00–05:00 KST)·Greengrass 검증","롤백 및 우선 배포 설정","펌웨어 테이블: S/N·모델·현재 버전·목표 버전·자동 업데이트","자동 업데이트 ON/OFF 필터"],
   ["OTA 작업 배포","자동 업데이트 필터링"],BLU),

  ("/admin/iot-auth","IoT 인증 관리","🔐",
   "AWS IoT Core Thing 구성과 디바이스 인증을 관리합니다. JITP 프로비저닝, 인증서, 정책 관리로 IoT 기기 보안을 강화합니다.",
   ["인증 설정: JITP·MQTT over TLS 1.3·Secrets Manager/KMS","IoT Thing 테이블: S/N·Thing명·인증서·정책·연결 상태","연결 상태 필터: Connected/Disconnected/Pending","S/N·Thing·인증서·정책 검색"],
   ["연결 상태별 필터링"],BLU),
]

# ══════════════════════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════════════════════
def main():
    prs=prs_new()
    cover(prs)
    toc(prs)
    overview(prs)

    # Section 1 장비 관리
    sec_slide(prs,"01","장비 관리","장비 라이프사이클  ·  장비 목록  ·  장비 이력",CYN)
    lifecycle(prs,**LC_EQUIP)
    for pg in PAGES:
        if pg[0] in ("/equipment","/equipment-logs"):
            feat(prs,pg[1],pg[2],pg[3],pg[4],pg[5],pg[0],pg[6])
            if pg[0]=="/equipment-logs" and SS_TABS.get("/equipment-logs"):
                tab_grid(prs,"장비 이력","📋",SS_TABS["/equipment-logs"],CYN)

    # Section 2 데이터 모니터링
    sec_slide(prs,"02","데이터 수집 / 모니터링","파이프라인  ·  메트릭 스트림",BLU)
    for pg in PAGES:
        if pg[0] in ("/data-pipeline",):
            feat(prs,pg[1],pg[2],pg[3],pg[4],pg[5],pg[0],pg[6])
    # 메트릭 스트림 접근 권한 없음 설명
    s=blank(prs,RGBColor(0x0D,0x14,0x23))
    R(s,0,0,Inches(7.6),SLIDE_H,CC); R(s,0,0,Inches(0.12),SLIDE_H,BLU)
    T(s,"📡",Inches(0.25),Inches(0.35),Inches(1.1),Inches(1.0),sz=42,align=PP_ALIGN.CENTER)
    T(s,"메트릭 스트림",Inches(1.4),Inches(0.45),Inches(6.0),Inches(0.85),sz=26,bold=True)
    R(s,Inches(0.3),Inches(1.45),Inches(7.0),Inches(0.04),BLU)
    T(s,"장비 센서 데이터를 실시간 스트리밍으로 확인하는 화면입니다.\nEvents · Diagnostics · Performance · Telemetry 탭으로 분류됩니다.",
      Inches(0.3),Inches(1.6),Inches(7.1),Inches(1.8),sz=14,col=LGT)
    T(s,"  •  실시간 Live/Paused 스트림 제어\n  •  탭별 메트릭 카운트 및 통계\n  •  장비 선택 및 S/N·메트릭·단위 검색\n  •  스트림 일시정지/재개",
      Inches(0.3),Inches(3.6),Inches(7.1),Inches(2.5),sz=13,col=LGT)
    T(s,"🔒  현재 데모 계정 권한으로는 접근이 제한됩니다.",
      Inches(0.3),Inches(9.0),Inches(7.1),Inches(0.8),sz=14,col=AMB)
    ss=SS.get("/metric-stream")
    if ss and os.path.exists(ss):
        rw=SLIDE_W-Inches(8.1); sh=rw*9/16; ry=(SLIDE_H-sh)/2
        s.shapes.add_picture(ss,Inches(7.9),ry,rw,sh)

    # Section 3 알람 관리
    sec_slide(prs,"03","알람 관리","알람 처리 흐름  ·  알람 목록  ·  알람 규칙",AMB)
    lifecycle(prs,**LC_ALARM)
    for pg in PAGES:
        if pg[0] in ("/alarms","/alarm-rules"):
            feat(prs,pg[1],pg[2],pg[3],pg[4],pg[5],pg[0],pg[6])

    # Section 4 서비스 관리
    sec_slide(prs,"04","서비스 관리","서비스 티켓 라이프사이클  ·  진행 현황  ·  A/S 기록",PRP)
    lifecycle(prs,**LC_SVC)
    for pg in PAGES:
        if pg[0] in ("/service-tickets","/service-progress","/as"):
            feat(prs,pg[1],pg[2],pg[3],pg[4],pg[5],pg[0],pg[6])

    # Section 5 현장 작업
    sec_slide(prs,"05","현장 작업","설치 / 점검 / 원격 제어",RGBColor(0x0F,0x76,0x6E))
    lifecycle(prs,**LC_FIELD)
    for pg in PAGES:
        if pg[0] in ("/inspection","/installation","/remote-control"):
            feat(prs,pg[1],pg[2],pg[3],pg[4],pg[5],pg[0],pg[6])

    # Section 6 부품 관리
    sec_slide(prs,"06","부품 관리","부품 조달 라이프사이클  ·  주문  ·  일정",GRN)
    lifecycle(prs,**LC_PARTS)
    for pg in PAGES:
        if pg[0] in ("/parts-orders","/parts-schedule"):
            feat(prs,pg[1],pg[2],pg[3],pg[4],pg[5],pg[0],pg[6])

    # Section 7 고객 / SLA
    sec_slide(prs,"07","고객 / SLA 관리","고객사 정보  ·  서비스 레벨 계약",RGBColor(0x0F,0x76,0x6E))
    for pg in PAGES:
        if pg[0] in ("/customers","/sla"):
            feat(prs,pg[1],pg[2],pg[3],pg[4],pg[5],pg[0],pg[6])

    # Section 8 보고서 / 설정
    sec_slide(prs,"08","보고서 / 시스템 설정","보고서  ·  알림 채널  ·  펌웨어 OTA  ·  IoT 인증",AMB)
    for pg in PAGES:
        if pg[0] in ("/reports","/settings/notifications","/settings/firmware","/admin/iot-auth"):
            feat(prs,pg[1],pg[2],pg[3],pg[4],pg[5],pg[0],pg[6])

    # Section 9 대시보드
    sec_slide(prs,"09","통합 대시보드","실시간 KPI  ·  플릿 현황판",BLU)
    for pg in PAGES:
        if pg[0]=="/dashboard":
            feat(prs,pg[1],pg[2],pg[3],pg[4],pg[5],pg[0],pg[6])

    prs.save(OUT_FILE)
    total=len(prs.slides)
    print(f"✅  사용자 설명서 저장 완료: {OUT_FILE}  ({total}슬라이드)")

main()
